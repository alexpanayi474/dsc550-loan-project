"""
Generate a single Post-Deployment Monitoring slide for the DSC550 presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── colour palette ──────────────────────────────────────────
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
RED    = RGBColor(0xE0, 0x4F, 0x5F)
BLUE   = RGBColor(0x3A, 0x7C, 0xBD)
GOLD   = RGBColor(0xD4, 0xA0, 0x1E)

# ── helpers ─────────────────────────────────────────────────
def add_bg(slide, color=DARK):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, size=18, bold=False,
                 color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    return txBox

def add_divider(slide, top, color=ACCENT, width=10):
    left = (13.333 - width) / 2
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_rounded_box(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ════════════════════════════════════════════════════════════
# SLIDE — POST-DEPLOYMENT MONITORING
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "09  Post-Deployment Monitoring", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

monitor_cards = [
    ("1. Actual Default Rates", RED,
     [
         "After deployment, the bank should track",
         "the actual default rate among approved",
         "loans and compare it against the ~3%",
         "base rate observed in the training data.",
         "",
         "If the realized rate deviates noticeably,",
         "it means the applicant population or the",
         "economic environment has changed, and the",
         "model's probability estimates — and therefore",
         "the optimal thresholds — may no longer",
         "be accurate.",
     ]),
    ("2. When to Recalibrate", GOLD,
     [
         "Recalibration should be considered when:",
         "",
         "• The realized default rate shifts",
         "  materially from the ~3% baseline",
         "",
         "• The cost structure used in the profit",
         "  calculation is updated",
         "",
         "• A significant macroeconomic event",
         "  (e.g. recession, policy change) alters",
         "  the lending environment",
         "",
         "In each case, re-run the profit-per-1,000",
         "sweep with updated data or costs.",
     ]),
    ("3. Sensitivity to Cost Changes", BLUE,
     [
         "The optimal thresholds depend directly on",
         "the costs assigned to each outcome (TN, FP,",
         "FN, TP) per stakeholder.",
         "",
         "↑ Regulatory penalties (CRO's FN cost)",
         "→ Threshold shifts lower, more rejections",
         "",
         "↑ Value of good customers (CGO's TN value)",
         "→ Threshold shifts higher, more approvals",
         "",
         "The threshold is not a fixed number but a",
         "policy lever that should be re-optimized",
         "whenever the cost assumptions are revised.",
     ]),
]

for i, (title, color, lines) in enumerate(monitor_cards):
    x = 0.4 + i * 4.2
    add_rounded_box(slide, x, 1.5, 3.9, 5.5, RGBColor(0x24, 0x36, 0x58))
    add_text_box(slide, x + 0.15, 1.6, 3.6, 0.5, title, 16, True, color)
    # colour divider under title
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x + 0.25), Inches(2.15), Inches(3.4), Pt(2))
    div.fill.solid()
    div.fill.fore_color.rgb = color
    div.line.fill.background()
    # body text
    txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.35), Inches(3.6), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_after = Pt(2)

# ── SAVE ────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "monitoring_slide.pptx")
prs.save(out_path)
print(f"Saved to: {out_path}")
