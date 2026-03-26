"""
Generate PowerPoint presentation for DSC550 Group Project:
Cyprus Community Bank - Loan Default Prediction
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

PLOTS = os.path.join(os.path.dirname(__file__), "plots")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── colour palette ──────────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x00, 0x00, 0x00)
DARK    = RGBColor(0x1B, 0x2A, 0x4A)   # dark navy
ACCENT  = RGBColor(0x2E, 0x86, 0xAB)   # teal
LIGHT   = RGBColor(0xF0, 0xF4, 0xF8)   # very light grey-blue
RED     = RGBColor(0xE0, 0x4F, 0x5F)
GREEN   = RGBColor(0x2D, 0x93, 0x6C)
BLUE    = RGBColor(0x3A, 0x7C, 0xBD)
GOLD    = RGBColor(0xD4, 0xA0, 0x1E)
GREY    = RGBColor(0x6B, 0x7B, 0x8D)

# ── helpers ─────────────────────────────────────────────────
def add_bg(slide, color=DARK):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def tx(shape, text, size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_text_box(slide, left, top, width, height, text, size=18, bold=False,
                 color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tx(txBox, text, size, bold, color, alignment, font_name)
    return txBox

def add_bullet_frame(slide, left, top, width, height, items, size=16, color=WHITE, spacing=Pt(8), font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox

def add_image_safe(slide, path, left, top, width=None, height=None):
    full = os.path.join(PLOTS, path)
    if os.path.exists(full):
        kwargs = {}
        if width:  kwargs["width"]  = Inches(width)
        if height: kwargs["height"] = Inches(height)
        slide.shapes.add_picture(full, Inches(left), Inches(top), **kwargs)
    else:
        add_text_box(slide, left, top, 4, 1, f"[Image not found: {path}]", 14, color=RED)

def add_divider(slide, top, color=ACCENT, width=10):
    left = (13.333 - width) / 2
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_rounded_box(slide, left, top, width, height, fill_color, text="", text_size=14, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = "Calibri"
    return shape

# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)
add_divider(slide, 2.6, ACCENT, 6)

add_text_box(slide, 1, 1.0, 11.3, 1.5,
             "Cyprus Community Bank", 44, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, 1, 1.8, 11.3, 1,
             "Loan Default Prediction", 32, False, ACCENT, PP_ALIGN.CENTER)
add_divider(slide, 2.7, ACCENT, 4)
add_text_box(slide, 1, 3.1, 11.3, 0.8,
             "DSC 550 — Group Project", 20, False, GREY, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.8, 11.3, 0.8,
             "Alexandros Panayi  ·  Maria Michaelidou  ·  Thanasis Kalos", 18, False, WHITE, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "Agenda", 36, True, WHITE, PP_ALIGN.LEFT)
add_divider(slide, 1.1, ACCENT, 11.5)

sections = [
    ("01", "The Problem", "Predicting loan defaults for a Cyprus community bank"),
    ("02", "Exploratory Data Analysis", "Understanding the data landscape & class imbalance"),
    ("03", "Model Training & Evaluation", "7 classifiers, SMOTE vs. unbalanced comparison"),
    ("04", "Model Selection & Calibration", "Choosing the best model for threshold optimization"),
    ("05", "Stakeholder Cost Structures", "Three executives, three different cost functions"),
    ("06", "Threshold Optimization", "From probabilities to business-optimal decisions"),
    ("07", "Stakeholder Tradeoff Analysis", "The 3×3 payoff matrix: who wins, who loses"),
    ("08", "CEO Decision Framework", "A range of defensible options — not one answer"),
    ("09", "Post-Deployment Monitoring", "Keeping the model and thresholds current"),
]

y = 1.5
for num, title, desc in sections:
    add_rounded_box(slide, 0.8, y, 0.7, 0.5, ACCENT, num, 16, WHITE)
    add_text_box(slide, 1.7, y - 0.03, 4, 0.5, title, 18, True, WHITE)
    add_text_box(slide, 5.5, y, 6.5, 0.5, desc, 14, False, GREY)
    y += 0.6

# ════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "01  The Problem", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_bullet_frame(slide, 0.8, 1.5, 11.5, 5.5, [
    "Cyprus Community Bank needs a data-driven system to decide which loan applications to approve or reject.",
    "A single wrong approval of a defaulting borrower can cost up to €45,000 in losses.",
    "A single wrong rejection of a good borrower sacrifices up to €10,000 in lost growth revenue.",
    "Three senior executives — CFO, CRO, CGO — each define \"optimal\" differently:",
    "    CFO (Maria Georgiou): Maximize profit — willing to accept moderate risk",
    "    CRO (Andreas Konstantinou): Minimize risk exposure — nearly zero tolerance for defaults",
    "    CGO (Elena Papadopoulou): Maximize growth — approves unless near-certain default",
    "",
    "Our task: Build a predictive model and find the optimal decision threshold",
    "for each stakeholder — then present the CEO with a range of defensible options.",
], size=17, color=WHITE, spacing=Pt(6))

# ════════════════════════════════════════════════════════════
# SLIDE 4 — EDA: CLASS IMBALANCE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "02  Exploratory Data Analysis", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_text_box(slide, 0.8, 1.3, 5.5, 0.5, "Severe Class Imbalance: ~3.1% Default Rate", 22, True, WHITE)
add_bullet_frame(slide, 0.8, 1.9, 5.5, 2.0, [
    "Only ~310 out of 10,000 loans default",
    "A naive \"approve all\" baseline achieves 96.9% accuracy",
    "→ Accuracy is a misleading metric here",
    "We must focus on AUC, Recall, and cost-sensitive metrics",
], size=15, color=WHITE, spacing=Pt(4))

add_image_safe(slide, "class_distribution.png", 6.8, 1.3, width=5.8)

# ════════════════════════════════════════════════════════════
# SLIDE 5 — EDA: FEATURE DISTRIBUTIONS
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "02  Feature Distributions by Default Status", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_image_safe(slide, "feature_distributions_continuous.png", 0.3, 1.3, width=12.7, height=5.5)

# ════════════════════════════════════════════════════════════
# SLIDE 6 — EDA: DISCRETE FEATURES
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "02  Discrete Feature Distributions", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_image_safe(slide, "feature_distributions_discrete.png", 0.3, 1.3, width=12.7, height=5.5)

# ════════════════════════════════════════════════════════════
# SLIDE 7 — EDA: CORRELATIONS & PURPOSE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "02  Key Predictors of Default", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_image_safe(slide, "correlation.png", 0.3, 1.3, width=6.0, height=4.5)
add_image_safe(slide, "default_by_purpose.png", 6.5, 1.3, width=6.2, height=4.0)

add_bullet_frame(slide, 6.8, 5.4, 5.5, 2.0, [
    "Strongest predictors: previous_defaults (+), debt_to_income (+)",
    "Protective: credit_score (−), age (−), income (−)",
    "Highest risk: debt_consolidation (4.2%), business (3.8%)",
], size=14, color=WHITE, spacing=Pt(3))

# ════════════════════════════════════════════════════════════
# SLIDE 8 — EDA SUMMARY
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "02  EDA — Key Takeaways", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_bullet_frame(slide, 0.8, 1.5, 11.5, 5.5, [
    "1. Severe Imbalance (~3.1% defaults) — naive baseline = 96.9% accuracy → misleading",
    "",
    "2. Continuous Features: Defaulters tend to be younger, lower income, lower credit score,",
    "    fewer years employed — indicating financially unstable profiles",
    "",
    "3. Discrete Features: previous_defaults is the single strongest signal;",
    "    more existing loans and longer loan terms also increase risk",
    "",
    "4. Correlation: previous_defaults and debt_to_income_ratio are the two strongest",
    "    positive predictors; credit_score is the strongest protective factor",
    "",
    "5. Loan Purpose: debt_consolidation (4.2%) and business (3.8%) carry the most risk;",
    "    home_improvement (1.9%) is safest",
], size=17, color=WHITE, spacing=Pt(3))

# ════════════════════════════════════════════════════════════
# SLIDE 9 — MODEL TRAINING
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "03  Model Training & Evaluation", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_text_box(slide, 0.8, 1.3, 11, 0.5, "7 classifiers trained on SMOTE-balanced data, evaluated on held-out test set", 18, False, GREY)

# Model comparison table
headers = ["Model", "Accuracy", "Precision", "Recall", "F1", "AUC"]
data = [
    ["Logistic Regression", "0.7880", "0.0923", "0.6613", "0.1621", "0.8161"],
    ["LDA",                 "0.7800", "0.0891", "0.6613", "0.1571", "0.8143"],
    ["Gradient Boosting",   "0.9570", "0.2273", "0.1613", "0.1887", "0.7610"],
    ["Random Forest",       "0.9620", "0.2308", "0.0968", "0.1364", "0.7415"],
    ["Naive Bayes",         "0.6735", "0.0583", "0.6290", "0.1067", "0.7337"],
    ["KNN",                 "0.8275", "0.0673", "0.3548", "0.1131", "0.6360"],
    ["Decision Tree",       "0.9295", "0.1376", "0.2419", "0.1754", "0.5967"],
]

table_l = 0.8
table_t = 2.0
table_w = 11.5
table_h = 4.0
rows_n = len(data) + 1
cols_n = len(headers)
table_shape = slide.shapes.add_table(rows_n, cols_n, Inches(table_l), Inches(table_t), Inches(table_w), Inches(table_h))
table = table_shape.table

# Style header
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.font.name = "Calibri"
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT

# Fill data
for i, row in enumerate(data):
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x24, 0x36, 0x58) if i % 2 == 0 else RGBColor(0x1B, 0x2A, 0x4A)
    # Highlight AUC column for top 2
    if i < 2:
        cell = table.cell(i + 1, 5)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2D, 0x93, 0x6C)

add_text_box(slide, 0.8, 6.2, 11, 0.5, "→ Top 2 by AUC: Logistic Regression (0.816) and LDA (0.814) — selected for tuning", 16, True, GREEN)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — BALANCED vs UNBALANCED
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "04  SMOTE vs. Unbalanced — Why Calibration Matters", 32, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_image_safe(slide, "calibration_comparison.png", 0.3, 1.3, width=7.0, height=3.2)

add_bullet_frame(slide, 7.5, 1.3, 5.3, 3.5, [
    "SMOTE inflates predicted probabilities",
    "→ shifts optimal F1 threshold from 0.13 to 0.86",
    "",
    "Unbalanced model keeps calibrated probs",
    "→ P(default) = 10% genuinely means 10% risk",
    "",
    "Since we optimize by sweeping thresholds,",
    "we DON'T need SMOTE to boost recall at 0.5",
], size=15, color=WHITE, spacing=Pt(4))

add_image_safe(slide, "f1_vs_threshold_balanced_comparison.png", 0.3, 4.5, width=7.0, height=2.8)

add_bullet_frame(slide, 7.5, 5.0, 5.3, 2.0, [
    "AUC comparison:",
    "  LR: Unbalanced 0.823 vs SMOTE 0.816",
    "  LDA: Unbalanced 0.810 vs SMOTE 0.814",
    "",
    "→ Proceed with UNBALANCED training",
], size=15, color=WHITE, spacing=Pt(3))

# ════════════════════════════════════════════════════════════
# SLIDE 11 — MODEL SELECTION & TUNING
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "04  Model Selection & Hyperparameter Tuning", 32, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_text_box(slide, 0.8, 1.3, 11, 0.5, "Why AUC-ROC as the selection metric?", 22, True, WHITE)
add_bullet_frame(slide, 0.8, 1.9, 5.5, 2.5, [
    "AUC measures ranking ability across ALL thresholds",
    "We need a model that separates defaults from non-defaults",
    "well regardless of which cutoff each stakeholder picks",
    "",
    "Unlike accuracy/F1, AUC is threshold-agnostic",
    "— perfect for our multi-stakeholder problem",
], size=15, color=WHITE, spacing=Pt(4))

add_text_box(slide, 0.8, 4.2, 11, 0.5, "Tuning Results (5-fold Stratified CV, scored by AUC)", 20, True, WHITE)

# Mini table for tuning
tune_headers = ["Model", "Best Params", "CV AUC", "Test AUC"]
tune_data = [
    ["Logistic Regression", "C = 1, penalty = L2", "0.8369", "0.8227"],
    ["LDA", "shrinkage = auto", "0.8303", "0.8104"],
]

tbl = slide.shapes.add_table(3, 4, Inches(0.8), Inches(4.8), Inches(11), Inches(1.5))
t = tbl.table
for j, h in enumerate(tune_headers):
    cell = t.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT

for i, row in enumerate(tune_data):
    for j, val in enumerate(row):
        cell = t.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x24, 0x36, 0x58)

add_text_box(slide, 0.8, 6.5, 11, 0.5, "→ Final model: Tuned Logistic Regression (AUC = 0.823) — trained on unbalanced data", 16, True, GREEN)

# ════════════════════════════════════════════════════════════
# SLIDE 12 — ROC CURVES
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "04  ROC Curves — All Models", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_image_safe(slide, "roc_all_models.png", 1.5, 1.3, width=10, height=5.8)

# ════════════════════════════════════════════════════════════
# SLIDE 13 — WHY LOGISTIC REGRESSION
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "04  Why Logistic Regression?", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

# Three cards
card_data = [
    ("Highest AUC (0.823)", "Best ability to separate defaults\nfrom non-defaults across all\npossible thresholds", GREEN),
    ("Calibrated Probabilities", "Predicted P(default) = 10%\ngenuinely means ~10% real risk\n→ thresholds are interpretable", BLUE),
    ("Full Transparency", "Odds ratios explain exactly\nwhy each loan is approved/rejected\n→ regulatory compliance advantage", GOLD),
]

for i, (title, desc, color) in enumerate(card_data):
    x = 0.8 + i * 4.1
    shape = add_rounded_box(slide, x, 1.5, 3.7, 3.5, RGBColor(0x24, 0x36, 0x58))
    add_text_box(slide, x + 0.2, 1.7, 3.3, 0.6, title, 20, True, color)
    add_text_box(slide, x + 0.2, 2.5, 3.3, 2.3, desc, 16, False, WHITE)

add_image_safe(slide, "roc_operating_points.png", 1.5, 5.2, width=4.5, height=2.2)

add_bullet_frame(slide, 6.5, 5.3, 6, 2.0, [
    "Unlike black-box models (Random Forest, GBM),",
    "Logistic Regression lets us explain every decision",
    "to regulators, customers, and the CEO",
], size=15, color=GREY, spacing=Pt(3))

# ════════════════════════════════════════════════════════════
# SLIDE 14 — STAKEHOLDER COST STRUCTURES (BUSINESS PERSPECTIVE)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "05  Stakeholder Cost Structures", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_text_box(slide, 0.8, 1.3, 11, 0.5, "Each executive values loan outcomes differently — this drives their preferred threshold", 17, False, GREY)

# Three stakeholder cards
sdata = [
    ("CFO — Maria Georgiou", "Maximize Profit", 
     "Approving a good loan: +€2,400\nApproving a default: −€15,000\nRejecting a good loan: €0\nRejecting a default: €0",
     "Break-even: 13.8%\n→ Reject only if >13.8% default risk",
     GREEN),
    ("CRO — Andreas Konstantinou", "Minimize Risk",
     "Approving a good loan: €0\nApproving a default: −€45,000\nRejecting a good loan: −€500\nRejecting a default: €0",
     "Break-even: 1.1%\n→ Reject if >1.1% default risk",
     RED),
    ("CGO — Elena Papadopoulou", "Maximize Growth",
     "Approving a good loan: +€8,000\nApproving a default: −€12,000\nRejecting a good loan: −€10,000\nRejecting a default: €0",
     "Break-even: 60.0%\n→ Reject only if >60% default risk",
     BLUE),
]

for i, (name, goal, costs, breakeven, color) in enumerate(sdata):
    x = 0.5 + i * 4.2
    shape = add_rounded_box(slide, x, 1.9, 3.9, 4.8, RGBColor(0x24, 0x36, 0x58))
    add_text_box(slide, x + 0.15, 2.0, 3.6, 0.5, name, 15, True, color)
    add_text_box(slide, x + 0.15, 2.5, 3.6, 0.4, goal, 13, False, GREY)
    add_text_box(slide, x + 0.15, 3.0, 3.6, 2.0, costs, 13, False, WHITE)
    add_rounded_box(slide, x + 0.3, 5.2, 3.3, 1.2, color, breakeven, 14, WHITE)

add_text_box(slide, 0.8, 6.9, 11, 0.5,
             "Ordering: CRO (1.1%) < CFO (13.8%) < CGO (60.0%) — from most conservative to most aggressive",
             15, True, GOLD, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 15 — THRESHOLD OPTIMIZATION
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "06  Threshold Optimization — Profit vs. Threshold", 32, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_image_safe(slide, "profit_vs_threshold.png", 0.3, 1.3, width=8.5, height=5.2)

add_bullet_frame(slide, 9.0, 1.5, 4.0, 5.0, [
    "Optimal Thresholds:",
    "",
    "CFO: t* = 0.14",
    "  → €1,906,800 profit/1k apps",
    "  → 95.4% approval rate",
    "",
    "CRO: t* = 0.01",
    "  → €−292,250 cost/1k apps",
    "  → 47.6% approval rate",
    "",
    "CGO: t* = 0.65",
    "  → €7,380,000 growth/1k",
    "  → 100% approval rate",
], size=14, color=WHITE, spacing=Pt(2))

# ════════════════════════════════════════════════════════════
# SLIDE 16 — ROC OPERATING POINTS
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "06  Stakeholder Operating Points on the ROC Curve", 32, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_image_safe(slide, "roc_operating_points.png", 1.5, 1.3, width=10, height=5.8)

# ════════════════════════════════════════════════════════════
# SLIDE 17 — TRADEOFF MATRIX
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "07  Stakeholder Tradeoff Analysis", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_text_box(slide, 0.8, 1.2, 11, 0.5,
             "If we adopt one executive's preferred threshold, what does each stakeholder gain or lose?", 17, False, GREY)

# 3x3 tradeoff table
trade_headers = ["Strategy", "CFO €/1k", "CRO €/1k", "CGO €/1k", "Approval %", "Sensitivity", "Specificity"]
trade_data = [
    ["CFO (t=0.14)",  "€1,906,800",  "−€1,008,500", "€6,822,000", "95.4%", "0.290", "0.962"],
    ["CRO (t=0.01)",  "€1,123,800",  "−€292,250",   "−€1,161,000", "47.6%", "0.968", "0.490"],
    ["CGO (t=0.65)",  "€1,860,600",  "−€1,395,000", "€7,380,000", "100%", "0.000", "1.000"],
]

tbl = slide.shapes.add_table(4, 7, Inches(0.5), Inches(1.8), Inches(12.3), Inches(2.5))
t = tbl.table
for j, h in enumerate(trade_headers):
    cell = t.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT

strategy_colors = [GREEN, RED, BLUE]
for i, row in enumerate(trade_data):
    for j, val in enumerate(row):
        cell = t.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x24, 0x36, 0x58) if i % 2 == 0 else RGBColor(0x1B, 0x2A, 0x4A)

# Key tradeoff statements
add_text_box(slide, 0.8, 4.6, 11, 0.5, "Key Tradeoffs for the CEO:", 20, True, GOLD)

add_bullet_frame(slide, 0.8, 5.1, 11.5, 2.5, [
    "CRO vs. CFO: Adopting the CRO's threshold costs the CFO €783,000/1k in profit,",
    "    but reduces risk exposure by €716,250/1k. Approval rate drops from 95.4% → 47.6%.",
    "",
    "CGO vs. CFO: Adopting the CGO's threshold gains €558,000/1k in growth value,",
    "    but zero defaults are caught (sensitivity = 0%). Risk exposure worsens by €386,500/1k.",
    "",
    "CGO vs. CRO (the extremes): Moving from CRO to CGO increases risk exposure by €1,102,750/1k.",
], size=14, color=WHITE, spacing=Pt(2))

# ════════════════════════════════════════════════════════════
# SLIDE 18 — CEO DECISION FRAMEWORK (OVERVIEW)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "08  CEO Decision Framework", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_text_box(slide, 0.8, 1.3, 11, 0.7,
             "There is no single \"right\" threshold. The CEO must choose within a range\nof defensible options based on the bank's strategic priorities.", 19, False, WHITE)

add_image_safe(slide, "ceo_framework.png", 0.3, 2.2, width=8.0, height=5.0)

add_bullet_frame(slide, 8.5, 2.5, 4.3, 4.5, [
    "The feasible range:",
    "  CRO extreme → t = 0.01",
    "  CFO optimum → t = 0.14",
    "  CGO extreme → t = 0.65",
    "",
    "Gold zone: compromise",
    "region between CRO & CFO",
    "where small profit sacrifices",
    "yield large risk reductions",
    "",
    "The threshold is NOT fixed —",
    "it's a dynamic control lever",
], size=14, color=WHITE, spacing=Pt(3))

# ════════════════════════════════════════════════════════════
# SLIDE 19 — CHEAP COMPROMISES
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "08  Finding the \"Cheap Compromises\"", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_text_box(slide, 0.8, 1.3, 11, 0.6,
             "Moving from the CFO's optimum (t=0.14) toward the CRO (t=0.01):", 18, False, GREY)

comp_h = ["Threshold", "CFO €/1k", "CRO €/1k", "CGO €/1k", "Approval %", "CFO Cost", "CRO Gain"]
comp_data = [
    ["0.138 (CFO)", "1,906,800", "−1,008,500", "6,822,000", "95.4%", "—",    "—"],
    ["0.120",       "1,899,900", "−989,000",   "6,720,000", "94.8%", "−€6,900",  "+€19,500"],
    ["0.101",       "1,862,700", "−996,750",   "6,441,000", "93.2%", "−€44,100", "+€11,750"],
    ["0.065",       "1,844,700", "−749,250",   "5,703,000", "88.1%", "−€62,100", "+€259,250"],
    ["0.047",       "1,761,900", "−682,750",   "4,881,000", "83.2%", "−€144,900","+€325,750"],
    ["0.028",       "1,604,100", "−506,250",   "3,195,000", "73.0%", "−€302,700","+€502,250"],
    ["0.010 (CRO)", "1,123,800", "−292,250",   "−1,161,000","47.6%", "−€783,000","+€716,250"],
]

tbl = slide.shapes.add_table(len(comp_data) + 1, 7, Inches(0.5), Inches(1.8), Inches(12.3), Inches(3.7))
t = tbl.table
for j, h in enumerate(comp_h):
    cell = t.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT

for i, row in enumerate(comp_data):
    for j, val in enumerate(row):
        cell = t.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        # Highlight first and last rows
        if i == 0:
            cell.fill.fore_color.rgb = RGBColor(0x2D, 0x70, 0x4F)
        elif i == len(comp_data) - 1:
            cell.fill.fore_color.rgb = RGBColor(0x8B, 0x30, 0x30)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x24, 0x36, 0x58)

add_text_box(slide, 0.8, 5.7, 11, 0.5, "Key Insight:", 20, True, GOLD)
add_bullet_frame(slide, 0.8, 6.2, 11.5, 1.2, [
    "Moving from t=0.14 to t=0.12 costs CFO only €6,900 but gains €19,500 in risk reduction",
    "→ This is a high-value, low-cost compromise the CEO should consider",
    "The marginal cost of further tightening accelerates — diminishing returns below t ≈ 0.05",
], size=15, color=WHITE, spacing=Pt(3))

# ════════════════════════════════════════════════════════════
# SLIDE 20 — WHAT DRIVES DEFAULT RISK
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "08  What Drives Default Risk — Model Interpretability", 30, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_image_safe(slide, "model_coefficients.png", 7.0, 1.3, width=5.8, height=4.2)

add_bullet_frame(slide, 0.8, 1.5, 6.0, 5.5, [
    "Top Risk Factors (increase default risk):",
    "  • avg_monthly_balance — strongest risk signal",
    "  • loan_term_months — longer terms = higher risk",
    "  • debt_to_income_ratio — moderate risk signal",
    "",
    "Top Protective Factors (decrease default risk):",
    "  • credit_score — strongest protection",
    "  • home_improvement purpose — lower risk",
    "  • education purpose — lower risk",
    "",
    "Why this matters for the CEO:",
    "  For CRO → risk flags are exactly what regulators watch",
    "  For CGO → protective factors identify safe customers to approve",
    "  For CFO → risk is concentrated in specific profiles,",
    "       enabling targeted rules rather than blanket rejection",
], size=15, color=WHITE, spacing=Pt(3))

# ════════════════════════════════════════════════════════════
# SLIDE 21 — POST-DEPLOYMENT MONITORING
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "09  Post-Deployment Monitoring & Threshold Management", 30, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_text_box(slide, 0.8, 1.2, 11, 0.6,
             "The threshold is a dynamic control lever — not a one-time decision", 20, True, WHITE)

# Four monitoring cards
monitor_cards = [
    ("Default Rate Tracking", RED,
     "Monitor monthly default rate vs. ~3% baseline\n\n"
     "If default rate > 3.5%:\n"
     "→ Shift threshold toward CRO (tighter)\n\n"
     "If default rate < 2.5%:\n"
     "→ Consider relaxing toward CFO/CGO"),
    ("Cost Sensitivity", GOLD,
     "Monitor changes in cost assumptions\n\n"
     "If regulatory penalties increase:\n"
     "→ Move threshold lower (more conservative)\n\n"
     "If customer lifetime value grows:\n"
     "→ Move threshold higher (more approvals)\n\n"
     "Re-evaluate quarterly"),
    ("Model Performance", BLUE,
     "Track AUC and feature distributions monthly\n\n"
     "If AUC drops below 0.75:\n"
     "→ Retrain model on fresh data\n\n"
     "Monitor credit_score & debt_to_income\n"
     "distributions for drift"),
    ("Threshold Recalibration", GREEN,
     "CEO treats threshold as movable within\n"
     "the feasible range (≈ 0.01 – 0.14)\n\n"
     "Adjust based on:\n"
     "  • Current risk tolerance\n"
     "  • Regulatory pressure\n"
     "  • Growth targets\n"
     "  • Macroeconomic conditions"),
]

for i, (title, color, desc) in enumerate(monitor_cards):
    x = 0.4 + i * 3.2
    shape = add_rounded_box(slide, x, 2.0, 3.0, 5.0, RGBColor(0x24, 0x36, 0x58))
    add_text_box(slide, x + 0.1, 2.1, 2.8, 0.5, title, 15, True, color)
    add_divider_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x + 0.2), Inches(2.6), Inches(2.6), Pt(2))
    add_divider_line.fill.solid()
    add_divider_line.fill.fore_color.rgb = color
    add_divider_line.line.fill.background()
    add_text_box(slide, x + 0.1, 2.75, 2.8, 4.0, desc, 12, False, WHITE)

# ════════════════════════════════════════════════════════════
# SLIDE 22 — TRIGGER RULES
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "09  Decision Trigger Rules", 34, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_text_box(slide, 0.8, 1.2, 11, 0.5, "Connecting metrics to concrete actions:", 18, False, GREY)

# Trigger rules table
trigger_h = ["Trigger Condition", "Action", "Direction"]
trigger_data = [
    ["Default rate rises above 3.5%",          "Tighten threshold",       "→ Move toward CRO (lower t)"],
    ["Default rate drops below 2.5%",           "Consider relaxing threshold", "→ Move toward CFO (higher t)"],
    ["Profit declining but risk stable",        "Re-optimize for profit",  "→ Move toward CFO optimum"],
    ["Loan approvals stagnating / growth down", "Relax threshold",         "→ Move toward CGO (higher t)"],
    ["Regulatory penalties increase by >20%",   "Shift threshold left",    "→ More conservative"],
    ["Customer lifetime value increases",       "Shift threshold right",   "→ More approvals"],
    ["Model AUC drops below 0.75",              "Retrain model entirely",  "→ New model + recalibrate"],
    ["Macro shift (recession, policy change)",  "Full threshold review",   "→ CEO convenes stakeholders"],
]

tbl = slide.shapes.add_table(len(trigger_data) + 1, 3, Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.0))
t = tbl.table
for j, h in enumerate(trigger_h):
    cell = t.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT

for i, row in enumerate(trigger_data):
    for j, val in enumerate(row):
        cell = t.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.name = "Calibri"
            p.alignment = PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x24, 0x36, 0x58) if i % 2 == 0 else RGBColor(0x1B, 0x2A, 0x4A)

# ════════════════════════════════════════════════════════════
# SLIDE 23 — SUMMARY / CONCLUSION
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, 0.8, 0.4, 11.5, 0.8, "Summary & Recommendation", 36, True, ACCENT)
add_divider(slide, 1.1, ACCENT, 11.5)

add_bullet_frame(slide, 0.8, 1.5, 11.5, 6.0, [
    "1. We built a transparent, calibrated Logistic Regression model (AUC = 0.823)",
    "    that separates defaulting from non-defaulting borrowers effectively.",
    "",
    "2. Three stakeholders have fundamentally different optimal thresholds:",
    "    • CRO: t = 0.01 (reject almost all risk → 47.6% approval, catches 96.8% of defaults)",
    "    • CFO: t = 0.14 (maximize profit → 95.4% approval, catches 29.0% of defaults)",
    "    • CGO: t = 0.65 (approve everyone → 100% approval, catches 0% of defaults)",
    "",
    "3. The threshold is not a fixed number — it is a dynamic control lever within a range",
    "    of defensible options (~0.01 to ~0.14). The CEO adjusts based on strategic priorities.",
    "",
    "4. \"Cheap compromises\" exist: moving from t=0.14 to t=0.12 costs the CFO only €6,900",
    "    but reduces risk exposure by €19,500 — high-value, low-cost middle ground.",
    "",
    "5. Post-deployment: monitor default rates, cost assumptions, and model performance.",
    "    Specific trigger rules (e.g., default > 3.5% → tighten, AUC < 0.75 → retrain)",
    "    keep the system adaptive and aligned with evolving bank strategy.",
], size=17, color=WHITE, spacing=Pt(3))

# ════════════════════════════════════════════════════════════
# SLIDE 24 — THANK YOU
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_divider(slide, 3.0, ACCENT, 6)
add_text_box(slide, 1, 2.0, 11.3, 1.5, "Thank You", 44, True, WHITE, PP_ALIGN.CENTER)
add_divider(slide, 3.0, ACCENT, 6)
add_text_box(slide, 1, 3.5, 11.3, 0.8, "Questions?", 28, False, GREY, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11.3, 0.8,
             "Alexandros Panayi  ·  Maria Michaelidou  ·  Thanasis Kalos",
             18, False, ACCENT, PP_ALIGN.CENTER)

# ── SAVE ────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "DSC550_Loan_Default_Presentation.pptx")
prs.save(out_path)
print(f"Presentation saved to: {out_path}")
print(f"Total slides: {len(prs.slides)}")
